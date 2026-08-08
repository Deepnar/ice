"""C12 D9: file → structured text blocks, one parser per format family.

A parser's job is narrow and mechanical: get the text out with whatever
structure the format actually carries (page numbers, headings, slide numbers,
sheet names), and refuse loudly when it cannot. Sectioning, chunking and every
downstream decision belong to `ingest.py`; nothing here knows about ICE.

Format policy (D9) — every parser is a pure-python wheel, because Track F's
end state is one installer and a system binary would fight it:

    pdf   pypdf          docx  python-docx      pptx  python-pptx
    xlsx  openpyxl       csv   pandas           html  beautifulsoup4
    text  read directly (.txt .md .rst .org + source files)

Deliberately unsupported, with an error naming what IS supported: legacy .doc,
.rtf, .odt, .epub, images. Obscure relative to the work.

Scanned PDFs are DETECTED AND REFUSED rather than ingested as an empty
document — a document that silently contains nothing is worse than one that
failed, because nothing downstream can tell the difference. OCR is deferred to
Track F as a pluggable extraction engine (see `extract.py`).
"""

import csv as csv_module
import io
import os
from dataclasses import dataclass, field
from typing import Optional

from src.api.config import settings
import structlog

logger = structlog.get_logger("ice.ingestion.documents.parsers")

# Table rendering caps (D11): a described table, not a row dump.


@dataclass
class Block:
    """One structural unit as the FILE defines it (a page, a heading's body, a
    slide, a sheet). Sections are derived from blocks by the chunker — a block
    may be split, but two blocks are never merged, so a section never straddles
    a page or heading boundary."""
    text: str
    meta: dict = field(default_factory=dict)


@dataclass
class DocumentParse:
    blocks: list
    file_type: str
    title: str
    page_count: Optional[int] = None


class UnsupportedDocument(ValueError):
    """The file type is not supported, or the file carries no extractable text."""


TEXT_EXTENSIONS = {
    ".txt", ".md", ".markdown", ".rst", ".org", ".text", ".log",
    ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".h", ".cpp", ".hpp",
    ".go", ".rs", ".rb", ".php", ".sh", ".bash", ".zsh", ".sql", ".r", ".m",
    ".swift", ".kt", ".scala", ".pl", ".lua", ".vim", ".el",
    ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".json", ".xml", ".css", ".scss", ".tex",
}
REFUSED_EXTENSIONS = {
    ".doc": "legacy Word (.doc)", ".rtf": "RTF", ".odt": "OpenDocument",
    ".epub": "EPUB", ".png": "images", ".jpg": "images", ".jpeg": "images",
    ".gif": "images", ".webp": "images", ".tiff": "images", ".bmp": "images",
}

SUPPORTED_SUMMARY = ("pdf, docx, pptx, xlsx, csv/tsv, html, and any plain-text "
                     "or source file (.txt .md .rst + code)")


def detect_file_type(path: str) -> str:
    """Extension → the parser family name. Raises UnsupportedDocument with a
    message naming what IS supported (an error that only says 'no' makes the
    user guess)."""
    ext = os.path.splitext(path)[1].lower()
    if ext in REFUSED_EXTENSIONS:
        raise UnsupportedDocument(
            f"{REFUSED_EXTENSIONS[ext]} is not supported — supported: "
            f"{SUPPORTED_SUMMARY}")
    if ext == ".pdf":
        return "pdf"
    if ext == ".docx":
        return "docx"
    if ext == ".pptx":
        return "pptx"
    if ext in (".xlsx", ".xlsm"):
        return "xlsx"
    if ext in (".csv", ".tsv"):
        return "csv"
    if ext in (".html", ".htm"):
        return "html"
    if ext in TEXT_EXTENSIONS:
        return ext.lstrip(".") or "txt"
    raise UnsupportedDocument(
        f"unsupported file type '{ext or '(none)'}' — supported: {SUPPORTED_SUMMARY}")


# ── per-format parsers ──────────────────────────────────────────────────────

def parse_text(path: str, file_type: str) -> DocumentParse:
    """Plain text and source files: one block. Markdown structure is not
    parsed here on purpose — the shared chunker already treats markdown
    headings as hard section boundaries (C3), so re-splitting would produce two
    disagreeing notions of a section."""
    with open(path, encoding="utf-8", errors="replace") as fh:
        raw = fh.read()
    if not raw.strip():
        raise UnsupportedDocument("file is empty")
    return DocumentParse(blocks=[Block(text=raw)], file_type=file_type,
                         title=os.path.basename(path))


def parse_pdf(path: str) -> DocumentParse:
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    try:
        reader = PdfReader(path)
    except PdfReadError as exc:
        raise UnsupportedDocument(f"cannot read PDF: {exc}")
    if getattr(reader, "is_encrypted", False):
        # Try the empty-password case pypdf can sometimes open, then give up:
        # a half-decrypted document would ingest as garbage.
        try:
            if reader.decrypt("") == 0:
                raise UnsupportedDocument(
                    "PDF is password-protected — decrypt it before uploading")
        except (NotImplementedError, PdfReadError) as exc:
            raise UnsupportedDocument(f"PDF is encrypted and cannot be read: {exc}")

    blocks, total_chars = [], 0
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception as exc:            # one bad page must not kill the file
            logger.warning("pdf_page_extract_failed", page=i + 1, error=str(exc))
            text = ""
        total_chars += len(text.strip())
        if text.strip():
            blocks.append(Block(text=text, meta={"page": i + 1}))

    n_pages = len(reader.pages)
    if total_chars == 0:
        # D9: refuse, loudly and specifically. This is the scanned-PDF case and
        # the user needs to know WHY nothing happened.
        raise UnsupportedDocument(
            f"no extractable text ({n_pages} pages, 0 characters) — this looks "
            "like a scan. OCR is not enabled; see the Track F extraction-engine "
            "item (Tika/Docling)")
    return DocumentParse(blocks=blocks, file_type="pdf",
                         title=os.path.basename(path), page_count=n_pages)


def parse_docx(path: str) -> DocumentParse:
    """Paragraphs grouped under the heading that precedes them, so a block is
    a section of the document as its author defined it."""
    import docx

    document = docx.Document(path)
    blocks: list = []
    current_heading, buffer = None, []

    def flush():
        body = "\n\n".join(p for p in buffer if p.strip())
        if body.strip():
            meta = {"heading": current_heading} if current_heading else {}
            blocks.append(Block(text=body, meta=meta))

    for para in document.paragraphs:
        style = (para.style.name if para.style is not None else "") or ""
        text = para.text
        if style.startswith("Heading") and text.strip():
            flush()
            current_heading, buffer = text.strip(), [f"## {text.strip()}"]
            continue
        buffer.append(text)
    flush()

    for t_index, table in enumerate(document.tables):
        rows = [" | ".join(cell.text.strip() for cell in row.cells)
                for row in table.rows]
        body = "\n".join(r for r in rows if r.strip(" |"))
        if body.strip():
            blocks.append(Block(text=body, meta={"table": t_index + 1}))

    if not blocks:
        raise UnsupportedDocument("no extractable text in the document")
    return DocumentParse(blocks=blocks, file_type="docx",
                         title=os.path.basename(path))


def parse_pptx(path: str) -> DocumentParse:
    """One block per slide: shape text plus the speaker notes, which are often
    where the actual argument lives."""
    from pptx import Presentation

    presentation = Presentation(path)
    blocks = []
    for i, slide in enumerate(presentation.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes:
            parts.append(f"Speaker notes: {notes}")
        body = "\n\n".join(parts)
        if body.strip():
            blocks.append(Block(text=body, meta={"slide": i + 1}))
    if not blocks:
        raise UnsupportedDocument("no extractable text in the presentation")
    return DocumentParse(blocks=blocks, file_type="pptx",
                         title=os.path.basename(path),
                         page_count=len(presentation.slides))


def parse_html(path: str) -> DocumentParse:
    from bs4 import BeautifulSoup

    with open(path, encoding="utf-8", errors="replace") as fh:
        soup = BeautifulSoup(fh.read(), "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    title = (soup.title.string.strip()
             if soup.title and soup.title.string else os.path.basename(path))
    text = soup.get_text("\n")
    text = "\n".join(line.strip() for line in text.splitlines() if line.strip())
    if not text.strip():
        raise UnsupportedDocument("no extractable text in the HTML")
    return DocumentParse(blocks=[Block(text=text)], file_type="html",
                         title=title)


# ── tabular (D11): a DESCRIBED table, never a row dump ──────────────────────

def _describe_frame(frame, name: str) -> str:
    """Schema + per-column statistics. This is what makes a spreadsheet
    *memory*: 'this file has a churn_rate column covering 2024, 12k rows' is a
    fact the graph can hold, where 12,000 row-chunks are noise."""
    lines = [f"Table: {name} ({len(frame):,} rows x {len(frame.columns)} columns)",
             "", "Columns:"]
    for col in frame.columns:
        series = frame[col]
        nulls = int(series.isna().sum())
        dtype = str(series.dtype)
        try:
            if str(series.dtype).startswith(("int", "float", "Int", "Float")):
                described = (f"numeric — min {series.min()}, max {series.max()}, "
                             f"mean {round(float(series.mean()), 4)}")
            elif "datetime" in dtype:
                described = f"dates — {series.min()} to {series.max()}"
            else:
                counts = series.astype("string").value_counts().head(settings.document_table_top_values)
                top = ", ".join(f"{v} ({c})" for v, c in counts.items())
                described = f"{series.nunique()} distinct — {top}"
        except Exception:                   # a column of mixed junk
            described = f"{dtype}"
        lines.append(f"- {col} ({dtype}): {described}; {nulls} null")
    return "\n".join(lines)


def _render_rows(frame, name: str) -> str:
    """Head rows plus a deterministic spread from the rest, as a markdown
    table. Deterministic because a re-ingest of the same file must produce the
    same sections — otherwise idempotency keys change and the document
    duplicates."""
    columns = list(frame.columns)[:settings.document_table_max_render_cols]
    head_n = min(settings.document_table_sample_rows // 2, len(frame))
    picked = list(range(head_n))
    remaining = len(frame) - head_n
    if remaining > 0:
        want = min(settings.document_table_sample_rows - head_n, remaining)
        step = max(1, remaining // want)
        picked += list(range(head_n, len(frame), step))[:want]

    lines = [f"Sample rows from {name} "
             f"({len(picked)} of {len(frame):,}):", "",
             "| " + " | ".join(str(c) for c in columns) + " |",
             "|" + "|".join(["---"] * len(columns)) + "|"]
    for i in picked:
        row = frame.iloc[i]
        lines.append("| " + " | ".join(
            str(row[c]).replace("|", "\\|")[:120] for c in columns) + " |")
    if len(frame.columns) > settings.document_table_max_render_cols:
        lines.append("")
        lines.append(f"({len(frame.columns) - settings.document_table_max_render_cols} further "
                     "columns are described above but not sampled here.)")
    return "\n".join(lines)


def _frame_blocks(frame, sheet: str, meta_key: str) -> list:
    return [Block(text=_describe_frame(frame, sheet), meta={meta_key: sheet,
                                                            "part": "schema"}),
            Block(text=_render_rows(frame, sheet), meta={meta_key: sheet,
                                                         "part": "sample"})]


def parse_csv(path: str) -> DocumentParse:
    import pandas as pd

    delimiter = "\t" if path.lower().endswith(".tsv") else None
    if delimiter is None:
        # Sniff, because a ';'-separated European CSV read as one column
        # produces a described table that is entirely wrong.
        with open(path, encoding="utf-8", errors="replace") as fh:
            head = fh.read(8192)
        try:
            delimiter = csv_module.Sniffer().sniff(head, ",;\t|").delimiter
        except csv_module.Error:
            delimiter = ","
    try:
        frame = pd.read_csv(path, sep=delimiter)
    except Exception as exc:
        raise UnsupportedDocument(f"cannot read the table: {exc}")
    if frame.empty and not len(frame.columns):
        raise UnsupportedDocument("the table is empty")
    name = os.path.basename(path)
    return DocumentParse(blocks=_frame_blocks(frame, name, "table"),
                         file_type="csv", title=name)


def parse_xlsx(path: str) -> DocumentParse:
    import pandas as pd

    try:
        sheets = pd.read_excel(path, sheet_name=None)
    except Exception as exc:
        raise UnsupportedDocument(f"cannot read the workbook: {exc}")
    blocks = []
    for sheet_name, frame in sheets.items():
        if frame.empty and not len(frame.columns):
            continue
        blocks.extend(_frame_blocks(frame, str(sheet_name), "sheet"))
    if not blocks:
        raise UnsupportedDocument("the workbook has no populated sheets")
    return DocumentParse(blocks=blocks, file_type="xlsx",
                         title=os.path.basename(path), page_count=len(sheets))


_PARSERS = {
    "pdf": parse_pdf, "docx": parse_docx, "pptx": parse_pptx,
    "html": parse_html, "csv": parse_csv, "xlsx": parse_xlsx,
}


def parse_file(path: str, file_type: Optional[str] = None) -> DocumentParse:
    """Dispatch to the family parser. `file_type` overrides the extension."""
    file_type = file_type or detect_file_type(path)
    parser = _PARSERS.get(file_type)
    if parser is not None:
        return parser(path)
    return parse_text(path, file_type)


def parse_blob(text: str, title: str, file_type: str = "txt") -> DocumentParse:
    """A pasted blob — no file, no structure beyond what the text carries."""
    if not text.strip():
        raise UnsupportedDocument("nothing to ingest")
    return DocumentParse(blocks=[Block(text=text)], file_type=file_type,
                         title=title)


def block_label(meta: dict) -> str:
    """The human-readable location of a block, for the provenance header."""
    if not meta:
        return ""
    if "page" in meta:
        return f"p.{meta['page']}"
    if "slide" in meta:
        return f"slide {meta['slide']}"
    if "heading" in meta:
        return f"§{meta['heading']}"
    if "sheet" in meta:
        return f"sheet {meta['sheet']}" + (
            f" ({meta['part']})" if meta.get("part") else "")
    if "table" in meta:
        part = meta.get("part")
        return f"table {meta['table']}" if not part else f"{meta['table']} ({part})"
    return ""


def read_bytes(path: str) -> bytes:
    with io.open(path, "rb") as fh:
        return fh.read()
